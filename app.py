import logging
import streamlit as st
import re
import time
import io
import tempfile
from pathlib import Path
from openai import OpenAI
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_VERTICAL_ANCHOR
from script import PowerPointGenerator, Config
from bs4 import BeautifulSoup
import markdown  
import tempfile  
from question_utils import generate_question_paper  
import pandas as pd
from auth import sign_up_user, verify_user, logout
import ssl
ssl._create_default_https_context = ssl._create_unverified_context

# ===== PROMPT TEMPLATE =====
PROMPT_TEMPLATE = """
Assume the role of an experienced {job_role} responsible for 
{expertise}. Using your real-world experience and responsibilities as context, you are tasked with creating a training presentation for fresh engineering graduates.

The presentation focuses on core Electronic Product practices related to the core skill: {core_skill}.

Each ChatGPT Canvas will focus on one micro-skill and can have 1 to 5 slides, depending on complexity.

Each slide must include:
- Title formatted as: # Slide X: Slide Title
- 3–5 bullet points (start each with '-'). Keep each under 15 words.
- If including a table, limit to max 5 columns & 6 rows.
- Add speaker notes at the end using: speaker notes:
- Make all content compatible with markdown.
- Leave space on the right for avatar — do not overcrowd.
- Avoid LaTeX and keep formatting clean.

Now generate the Canvas (1 to 5 slides starting from Slide {slide_start}) for the micro-skill:
{microskill_title} — {microskill_details}
"""

# ===== STREAMLIT APP CONFIGURATION =====
st.set_page_config(
    page_title="Training PPT Generator",
    page_icon="📊",
    layout="wide",
    initial_sidebar_state="expanded"
)

# ===== HELPER FUNCTIONS =====
@st.cache_data
def query_openai(prompt, api_key):
    """Query OpenAI API with caching"""
    try:
        client = OpenAI(api_key=api_key)
        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {"role": "system", "content": "You are a markdown slide generator."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.5,
            max_tokens=3000
        )
        return response.choices[0].message.content, None
    except Exception as e:
        return None, str(e)

def parse_table(lines):
    """Parse markdown table lines into data structure"""
    table_data = []
    for line in lines:
        if line.strip().startswith("|") and "|" in line and not set(line.strip()).issubset({"|", "-", " ", ":"}):
            row = [cell.strip() for cell in line.strip().strip("|").split("|")]
            table_data.append(row)
    return table_data

# def markdown_to_ppt(markdown):
#     """Convert markdown to PowerPoint presentation"""
#     prs = Presentation()
#     prs.slide_width = Inches(10)
#     prs.slide_height = Inches(7.5)
#     notes_out = []

#     slide_chunks = re.findall(r"# Slide \d+: .*?(?=(?:# Slide \d+:|\Z))", markdown, re.DOTALL)

#     for idx, chunk in enumerate(slide_chunks):
#         lines = chunk.strip().splitlines()
#         title_line = lines[0].strip()
#         title = title_line.split(":", 1)[-1].strip()

#         bullet_lines, speaker_notes, table_lines = [], "", []
#         collecting_notes = False

#         for line in lines[1:]:
#             if line.strip().lower().startswith("speaker notes:"):
#                 collecting_notes = True
#                 speaker_notes = line.split(":", 1)[-1].strip()
#             elif collecting_notes:
#                 speaker_notes += " " + line.strip()
#             elif re.match(r"^[-*•]\s+", line.strip()):
#                 bullet_lines.append(re.sub(r"^[-*•]\s+", "", line.strip()))
#             elif line.strip().startswith("|"):
#                 table_lines.append(line)

#         slide = prs.slides.add_slide(prs.slide_layouts[5])
#         title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
#         tf = title_box.text_frame
#         tf.text = title
#         tf.paragraphs[0].font.size = Pt(28)
#         tf.paragraphs[0].font.bold = True
#         tf.paragraphs[0].alignment = PP_ALIGN.CENTER

#         if table_lines:
#             table_data = parse_table(table_lines)
#             if table_data and len(table_data[0]) <= 5 and len(table_data) <= 6:
#                 table_shape = slide.shapes.add_table(len(table_data), len(table_data[0]), Inches(1), Inches(1.5), Inches(8), Inches(3)).table
#                 for i, row in enumerate(table_data):
#                     for j, cell_text in enumerate(row):
#                         cell = table_shape.cell(i, j)
#                         cell.text = cell_text
#                         cell.text_frame.paragraphs[0].font.size = Pt(16)

#         if bullet_lines:
#             bullet_box = slide.shapes.add_textbox(Inches(0.63), Inches(2.25), Inches(7.5), Inches(3))
#             tf_bullets = bullet_box.text_frame
#             tf_bullets.clear()
#             tf_bullets.word_wrap = True
#             tf_bullets.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
#             bullet_box.text_frame.auto_size = True
        
#             # Set margins
#             bullet_box.text_frame.margin_left = Inches(0.098)
#             bullet_box.text_frame.margin_right = Inches(0.098)
#             bullet_box.text_frame.margin_top = Inches(0.051)
        
#             for i, bullet in enumerate(bullet_lines[:5]):
#                 p = tf_bullets.paragraphs[0] if i == 0 else tf_bullets.add_paragraph()
#                 p.text = bullet
#                 p.level = 0
#                 p.font.size = Pt(20)
#                 p.font.name = 'Calibri'
#                 p.alignment = PP_ALIGN.LEFT

#         slide.notes_slide.notes_text_frame.text = speaker_notes or "(No speaker notes provided.)"
#         notes_out.append(f"Slide {idx + 1} - {title}\n{speaker_notes or '(No speaker notes provided.)'}\n\n")

#     return prs, notes_out

def generate_ppt_files(job_role, expertise, core_skill, microskills_text, api_key, progress_callback=None):
    """Generate PPT and notes files from inputs"""
    all_markdown = []
    slide_count = 1
    total_lines = len([line for line in microskills_text.strip().split('\n') if '|' in line])
    
    for i, line in enumerate(microskills_text.strip().split('\n')):
        if "|" not in line:
            continue
        
        if progress_callback:
            progress_callback(f"Processing micro-skill {i+1} of {total_lines}...", (i+1) / total_lines)
        
        title, details = line.strip().split("|", 1)
        prompt = PROMPT_TEMPLATE.format(
            job_role=job_role,
            expertise=expertise,
            core_skill=core_skill,
            microskill_title=title.strip(),
            microskill_details=details.strip(),
            slide_start=slide_count
        )
        
        markdown_content, error = query_openai(prompt, api_key)
        if error:
            return None, None, None, error

        slide_count += markdown_content.count("# Slide")
        all_markdown.append(markdown_content)
        time.sleep(0.5)  # Rate limiting

    full_markdown = "\n\n".join(all_markdown)

    if progress_callback:
        progress_callback("Adding structure and personalisation...", 0.85)
    # Convert markdown to HTML
    html = markdown.markdown(full_markdown, extensions=['tables'])

    # Parse HTML using BeautifulSoup
    soup = BeautifulSoup(html, "lxml")
    content_div = soup.body  # The content to pass into PowerPointGenerator

    if progress_callback:
        progress_callback("Converting to PowerPoint...", 0.9)

    # Initialize generator with your existing style logic
    generator = PowerPointGenerator(Config(), logging.getLogger("pptgen"))

    # Generate PowerPoint into a temp path
    ppt_path = Path(tempfile.mktemp(suffix=".pptx"))
    success = generator.create_enhanced_presentation(content_div, ppt_path, title=core_skill)

    if not success:
        return None, None, None, "PowerPoint generation failed."

    prs = Presentation(ppt_path)
    notes_out = [f"Slide {i}: {n}" for i, n in generator.speaker_notes_txt]
    
        
    return prs, full_markdown, notes_out, None

def show_auth_ui():
    st.markdown("## ")  # vertical spacing

    # Center the full wrapper using columns
    col1, col2, col3 = st.columns([1, 2.5, 1])
    with  col2:
        with st.container():
            # Open wrapper box
            

            # --- All form content inside this box ---
            from pathlib import Path
            watermark_path = str(Path(__file__).parent / "assets" / "watermark.jpg")
            st.markdown("<div style='display:flex; justify-content:center; align-items:center; margin-bottom:10px;'>", unsafe_allow_html=True)
            st.image(watermark_path, width=180)
            st.markdown("</div>", unsafe_allow_html=True)
            st.markdown("<h2 style='text-align:center; color:#2b7cff; margin-bottom:0;'>Training PPT Generator</h2>", unsafe_allow_html=True)
            st.markdown("<p style='text-align:center; font-size:18px; color:#2b7cff;'>Please log in or register</p>", unsafe_allow_html=True)

            mode = st.radio("Choose Action", ["Login", "Sign Up"], horizontal=True)

            username = st.text_input("Username")
            password = st.text_input("Password", type="password")

            if mode == "Sign Up":
                email = st.text_input("Email")
                if st.button("Register"):
                    success, msg = sign_up_user(username, email, password)
                    if success:
                        st.success("Registration successful")
                    else:
                        st.error(msg)


            elif mode == "Login":
                if st.button("Login"):
                    user = verify_user(username, password)
                    if user:
                        st.session_state["user"] = user["username"]
                        st.success(f"Welcome, {user['username']}!")
                        st.rerun()
                    else:
                        st.error("Invalid credentials")

            # Close the wrapper
            st.markdown("</div>", unsafe_allow_html=True)



# ===== MAIN STREAMLIT APP =====
def main():
    if "user" not in st.session_state:
        show_auth_ui()
        st.stop()
    else:
        st.sidebar.success(f"Logged in as: {st.session_state['user']}")
        if st.sidebar.button("Logout"):
            logout()

    st.title("📊 Training PowerPoint Generator")
    st.markdown("Generate professional training presentations using AI")

    # Sidebar for configuration
    with st.sidebar:
        st.header("🔑 API Configuration")
        api_key = st.text_input("OpenAI API Key", type="password", help="Enter your OpenAI API key")
        
        st.header("📋 Training Details")
        job_role_options = [
            "Product Designer",
            "Domain Expert and V&V Engineer",
            "PCB Design Engineer",
            "Electronic Product Integration Engineer",
            "Procurement Specialist",
            "Mechanical Designer for Electro/Electronic Product",
            "Product Manager (Techno-Managerial)",
            "Firmware / Software Developer",
            "other"
        ]
        job_role_selected = st.selectbox(
            "Job Role",
            job_role_options,
            index=0
        )
        if job_role_selected == "other":
            job_role_custom = st.text_input("Please specify your job role", key="custom_job_role")
            job_role = job_role_custom.strip() if job_role_custom.strip() else "other"
        else:
            job_role = job_role_selected
        expertise = st.text_area("Required Expertise", placeholder="e.g., circuit design and testing")
        core_skill = st.text_input("Core Skill", placeholder="e.g., PCB Design")

    # Main content area
    col1, col2 = st.columns([2, 1])
    
    with col1:
        st.header("🎯 Micro-Skills Input")
        st.markdown("Enter micro-skills in the format: **Title|Details** (one per line)")
        
        microskills_text = st.text_area(
            "Micro-Skills",
            height=300,
            placeholder="Example:\nSoldering Basics|Understanding different soldering techniques and tools\nCircuit Analysis|Analyzing basic electronic circuits using Ohm's law",
            help="Each line should contain: Title|Description"
        )
        
        # Preview parsed micro-skills
        if microskills_text:
            st.subheader("📝 Parsed Micro-Skills Preview")
            parsed_skills = []
            for line in microskills_text.strip().split('\n'):
                if '|' in line:
                    title, details = line.strip().split('|', 1)
                    parsed_skills.append({"Title": title.strip(), "Details": details.strip()})
            
            if parsed_skills:
                st.dataframe(parsed_skills, use_container_width=True)
            else:
                st.warning("No valid micro-skills found. Please use the format: Title|Details")

    with col2:
        st.header("⚙️ Generation Settings")
        
        # Validation
        is_valid = all([api_key, job_role, expertise, core_skill, microskills_text])

        # Generation options
        st.markdown("**Select what you want to generate:**")
        col_gen1, col_gen2 = st.columns(2)
        with col_gen1:
            generate_ppt = st.checkbox("Generate PPT", value=True, key="generate_ppt")
        with col_gen2:
            generate_qp = st.checkbox("Generate Question Paper", value=False, key="generate_qp")

        # Unified generate button
        generate_btn = st.button(
            "🚀 Generate",
            disabled=not is_valid or (not generate_ppt and not generate_qp),
            use_container_width=True,
            type="primary"
        )

        # Show missing fields error ONLY after clicking Generate
        if generate_btn and not is_valid:
            missing_fields = []
            if not api_key: missing_fields.append("API Key")
            if not job_role: missing_fields.append("Job Role")
            if not expertise: missing_fields.append("Expertise")
            if not core_skill: missing_fields.append("Core Skill")
            if not microskills_text: missing_fields.append("Micro-Skills")
            st.error(f"Missing required fields: {', '.join(missing_fields)}")

    # Generation process
    # --- Persist generated files in session_state for persistent download buttons ---
    if 'ppt_buffer' not in st.session_state:
        st.session_state['ppt_buffer'] = None
    if 'notes_content' not in st.session_state:
        st.session_state['notes_content'] = None
    if 'full_markdown' not in st.session_state:
        st.session_state['full_markdown'] = None
    if 'last_core_skill' not in st.session_state:
        st.session_state['last_core_skill'] = None
    if 'question_excel_buffer' not in st.session_state:
        st.session_state['question_excel_buffer'] = None

    if generate_btn and is_valid:
        # Progress tracking
        progress_bar = st.progress(0)
        status_text = st.empty()

        def update_progress(message, progress):
            status_text.text(message)
            progress_bar.progress(progress)

        try:
            if generate_ppt:
                with st.spinner("🔄 Generating your presentation..."):
                    prs, full_markdown, notes_out, error = generate_ppt_files(
                        job_role, expertise, core_skill, microskills_text, api_key, update_progress
                    )

                    if error:
                        st.error(f"❌ Error generating presentation: {error}")
                        return

                    update_progress("Finalizing files...", 1.0)

                    # Create downloadable files and store in session_state
                    ppt_buffer = io.BytesIO()
                    prs.save(ppt_buffer)
                    ppt_buffer.seek(0)
                    notes_content = ''.join(notes_out)

                    st.session_state['ppt_buffer'] = ppt_buffer.getvalue()
                    st.session_state['notes_content'] = notes_content
                    st.session_state['full_markdown'] = full_markdown
                    st.session_state['last_core_skill'] = core_skill

                    # Success message
                    st.success("✅ Presentation generated successfully!")

                    # Clear progress indicators
                    progress_bar.empty()
                    status_text.empty()

            if generate_qp:
                with st.spinner("📝 Generating your question paper..."):
                    df, error = generate_question_paper(microskills_text, query_openai, api_key)
                    if error:
                        st.error(f"❌ {error}")
                        st.session_state['question_excel_buffer'] = None
                    else:
                        st.success("✅ Question paper generated successfully!")
                        excel_buffer = io.BytesIO()
                        with pd.ExcelWriter(excel_buffer, engine='xlsxwriter') as writer:
                            df.to_excel(writer, sheet_name='Questions', index=False)
                            workbook = writer.book
                            worksheet = writer.sheets['Questions']
                            wrap_format = workbook.add_format({'text_wrap': True})
                            worksheet.set_column('A:Z', 25, wrap_format)
                        excel_buffer.seek(0)
                        st.session_state['question_excel_buffer'] = excel_buffer.getvalue()
                        st.session_state['last_core_skill'] = core_skill
                        with st.expander("📋 Preview Question Table"):
                            st.dataframe(df)


        except Exception as e:
            st.error(f"❌ An unexpected error occurred: {str(e)}")
            progress_bar.empty()
            status_text.empty()


    # --- Show download buttons if files are available in session_state ---
    any_download = (
        st.session_state.get('ppt_buffer') and st.session_state.get('notes_content') and st.session_state.get('full_markdown')
    ) or st.session_state.get('question_excel_buffer')
    if any_download:
        cols = st.columns(4)
        core_skill_for_file = st.session_state.get('last_core_skill', 'presentation')
        col_idx = 0
        if st.session_state.get('ppt_buffer') and st.session_state.get('notes_content') and st.session_state.get('full_markdown'):
            with cols[col_idx]:
                st.download_button(
                    label="📄 Download PowerPoint",
                    data=st.session_state['ppt_buffer'],
                    file_name=f"{core_skill_for_file.replace(' ', '_')}_training.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )
            col_idx += 1
            with cols[col_idx]:
                st.download_button(
                    label="📝 Download Speaker Notes",
                    data=st.session_state['notes_content'],
                    file_name=f"{core_skill_for_file.replace(' ', '_')}_notes.txt",
                    mime="text/plain"
                )
            col_idx += 1
            with cols[col_idx]:
                st.download_button(
                    label="📋 Download Markdown",
                    data=st.session_state['full_markdown'],
                    file_name=f"{core_skill_for_file.replace(' ', '_')}_canvas.md",
                    mime="text/markdown"
                )
            col_idx += 1
        if st.session_state.get('question_excel_buffer'):
            with cols[col_idx]:
                st.download_button(
                    label="📥 Download Question Paper (Excel)",
                    data=st.session_state['question_excel_buffer'],
                    file_name=f"{core_skill_for_file.replace(' ', '_')}_questions.xlsx",
                    mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
                )

        # Show preview of generated content (if available)
        if st.session_state.get('full_markdown'):
            with st.expander("🔍 Preview Generated Content"):
                st.subheader("Generated Slides (Markdown)")
                st.code(st.session_state['full_markdown'][:2000] + "..." if len(st.session_state['full_markdown']) > 2000 else st.session_state['full_markdown'], language="markdown")
                st.subheader("Speaker Notes Preview")
                st.text(st.session_state['notes_content'][:1000] + "..." if len(st.session_state['notes_content']) > 1000 else st.session_state['notes_content'])

    # Footer
    st.markdown("---")
    st.markdown(
        """
        <div style='text-align: center; color: #666;'>
            <p><small>Ensure your API key has sufficient credits for generation</small></p>
        </div>
        """,
        unsafe_allow_html=True
    )

if __name__ == "__main__":
    main()